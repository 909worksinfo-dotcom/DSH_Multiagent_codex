from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import re

ROOT = Path(__file__).parent
OUT = ROOT / "RSI递归自我提升_硅谷前沿与投资研究_30页.pptx"
W, H = 13.333333, 7.5
COLORS = {
    "ink": "08111F", "slate": "132238", "slate2": "1B304C", "ice": "DDEBFA",
    "muted": "8FA7C2", "cyan": "49D6C9", "amber": "FFB454", "risk": "FF6B6B",
    "white": "F6FBFF", "line": "29415F", "black": "020812"
}

def rgb(key):
    s = COLORS.get(key, key).replace("#", "")
    return RGBColor.from_string(s)

def add_rect(slide, x, y, w, h, fill="slate", line="line", radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line); shp.line.width = Pt(0.8)
    return shp

def add_text(slide, text, x, y, w, h, size=18, color="ice", bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin); tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box

def add_rich(slide, title, body, x, y, w, h, accent="cyan"):
    add_rect(slide, x, y, w, h, "slate", "line", True)
    add_text(slide, title, x+0.18, y+0.16, w-0.36, 0.38, 17, accent, True)
    add_text(slide, body, x+0.18, y+0.58, w-0.36, h-0.72, 13.2, "ice")

def header(slide, num, kicker, title):
    add_text(slide, kicker, 0.64, 0.33, 10.8, 0.25, 10, "cyan", True, "Courier New")
    add_text(slide, title, 0.64, 0.67, 11.6, 0.62, 26, "white", True)
    add_text(slide, f"{num:02d}", 12.1, 0.38, 0.55, 0.3, 11, "muted", False, "Courier New", PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.37), Inches(12.05), Inches(0.008))
    line.fill.solid(); line.fill.fore_color.rgb = rgb("line"); line.line.fill.background()

def circle_loop(slide, cx=3.35, cy=4.25, r=1.65, center="RSI", nodes=None):
    nodes = nodes or ["PROPOSE", "EVALUATE", "SELECT", "DEPLOY"]
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    c.fill.background(); c.line.color.rgb = rgb("line"); c.line.width = Pt(1.5)
    positions = [(cx-0.65,cy-r-0.22),(cx+r-0.5,cy-0.23),(cx-0.65,cy+r-0.22),(cx-r-0.8,cy-0.23)]
    for text,(x,y) in zip(nodes,positions):
        add_rect(slide,x,y,1.3,0.46,"slate","cyan",True)
        add_text(slide,text,x+0.03,y+0.11,1.24,0.24,10,"cyan",True,"Courier New",PP_ALIGN.CENTER)
    add_text(slide,center,cx-0.75,cy-0.28,1.5,0.56,28,"cyan",True,"Aptos Display",PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)

def clean_text(tag):
    if not tag: return ""
    return re.sub(r"\s+", " ", tag.get_text(" ", strip=True)).strip()

html = BeautifulSoup((ROOT / "index.html").read_text(encoding="utf-8"), "html.parser")
sections = html.select("deck-stage > section")
assert len(sections) == 30
prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

for idx, sec in enumerate(sections, 1):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb("ink")
    texts = [clean_text(x) for x in sec.find_all(["h1","h2","h3","p","li","td","th","span","b","strong"], recursive=True)]
    title_tag = sec.find("h1") or sec.find("h2")
    title = clean_text(title_tag)
    kicker_tag = sec.select_one(".kicker")
    kicker = clean_text(kicker_tag) or "RECURSIVE SELF-IMPROVEMENT"

    if idx == 1:
        # cover
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(H))
        band.fill.solid(); band.fill.fore_color.rgb = rgb("cyan"); band.line.fill.background()
        add_text(slide, "RECURSIVE SELF-IMPROVEMENT · INVESTMENT RESEARCH", 0.72, 0.58, 7.2, 0.3, 11, "cyan", True, "Courier New")
        add_text(slide, "RSI：从模型升级到\n改进改进过程", 0.72, 1.32, 7.6, 1.65, 38, "white", True, "Aptos Display")
        add_text(slide, "硅谷前沿研究、技术路径与全球产业地图\n面向 VC 与 AI 行业从业者", 0.75, 3.27, 6.9, 0.85, 19, "muted")
        add_text(slide, "2026.08 · 30 SLIDES · RESEARCH EDITION", 0.75, 5.98, 6.4, 0.3, 11, "cyan", True, "Courier New")
        circle_loop(slide, 10.35, 3.7, 1.75)
        continue

    header(slide, idx, kicker, title)

    # specialized layouts
    if idx == 2:
        cards = [("0→2Y","可验证垂直闭环","代码、数学、芯片、实验设计等拥有低成本 verifier 的领域最先产生 ROI。\n\n投资：Agent evaluation、sandbox、trajectory data。","cyan"),("2→5Y","跨任务 Agentic RSI","系统重写工作流、工具与记忆，并把改进迁移到相邻任务。\n\n门槛：信用分配、隐藏评测、可回滚部署。","amber"),("5Y+","元级递归","系统改写学习算法、训练配方甚至架构；潜力最大，当前缺乏单调提升保证。","risk")]
        for i,(m,t,b,c) in enumerate(cards):
            x=0.65+i*4.12; add_rect(slide,x,1.75,3.82,4.85,"slate","line",True); add_text(slide,m,x+0.22,2.0,3.3,0.65,31,c,True,"Aptos Display"); add_text(slide,t,x+0.22,2.75,3.25,0.45,18,"white",True); add_text(slide,b,x+0.22,3.38,3.3,2.6,14,"ice")
    elif idx == 3:
        add_text(slide,"系统利用自身产出，改进产生下一次改进的机制；改进后的系统再重复该过程。",0.72,1.85,5.3,1.6,25,"white",True)
        add_text(slide,"RSI = self-reference + closed-loop evaluation + retained change + repeated iteration",0.72,4.0,5.2,0.95,13,"cyan",True,"Courier New")
        rows=[("Continual Learning","部分","吸收新数据不等于改进学习机制"),("AutoML / NAS","部分","外部固定搜索器通常不是自我改进"),("Self-training","条件成立","需形成可持续、可保留的能力闭环"),("Agent reflection","通常否","一次反思但不改变未来系统"),("Self-modifying agent","强相关","改变可执行策略且经验证保留")]
        y=1.78
        for a,b,c in rows:
            add_rect(slide,6.35,y,6.25,0.8,"slate","line"); add_text(slide,a,6.55,y+0.16,1.85,0.35,13,"white",True); add_text(slide,b,8.45,y+0.16,1.1,0.35,12,"cyan",True); add_text(slide,c,9.55,y+0.12,2.8,0.47,12.3,"ice"); y+=0.9
    elif idx in (4,8,30):
        center = "Δ > 0" if idx==4 else ("DATA" if idx==8 else "INFO")
        nodes = ["PROPOSE","EXECUTE","EVALUATE","RETAIN"] if idx==4 else (["GENERATE","CRITIQUE","TRAIN","REFRESH"] if idx==8 else ["TASK","SAMPLE","VERIFY","LEARN"])
        circle_loop(slide,3.3,4.25,1.65,center,nodes)
        cards = []
        if idx==4: cards=[("Generator","产生 prompt、数据、代码、架构或训练策略候选"),("Executor","隔离环境运行候选，记录完整 trajectory"),("Verifier","隐藏任务、单元测试或真实 KPI 判断净提升"),("Selector","综合性能、成本、安全后接受或回滚")]
        elif idx==8: cards=[("代表范式","Self-Instruct、STaR、Constitutional AI、Self-Rewarding LMs、RLAIF。"),("成立条件","生成器与验证器存在能力差异；规则、工具或真实反馈提供增量信息。"),("主要风险","错误放大、模式收缩、数据血缘不清与评判器偏差。")]
        else: cards=[("可靠信息源","代码执行、证明检查、模拟器、检索事实、人类偏好、真实业务结果。"),("防坍缩设计","保留真实数据锚点；追踪血缘；按不确定性采样；周期性外部校准。"),("投资判断","价值在高熵任务生成 + 低偏差验证 + 训练归因的完整供应链。")]
        y=1.75
        for t,b in cards:
            add_rich(slide,t,b,6.2,y,6.25,1.1); y+=1.25
    elif idx == 5:
        levels=[("L0","固定模型\n人工升级"),("L1","会话反思\n不持久化"),("L2","持久化 prompt / memory / workflow"),("L3","改写代码、工具与数据管线"),("L4","改写训练配方、学习器与架构"),("L5","跨域开放式递归")]
        for i,(l,b) in enumerate(levels):
            x=.65+i*2.03; ht=1.25+i*.6; y=6.45-ht; add_rect(slide,x,y,1.83,ht,"slate","cyan" if i<5 else "amber"); add_text(slide,l,x+.14,y+.16,1.45,.42,22,"white",True); add_text(slide,b,x+.14,y+.7,1.5,ht-.82,11.8,"ice")
    elif idx in (6,14,25,27):
        # fit HTML table into PPT table
        table = sec.find("table"); rows=table.find_all("tr"); data=[[clean_text(c) for c in r.find_all(["th","td"])] for r in rows]
        nrows,ncols=len(data),max(map(len,data)); shape=slide.shapes.add_table(nrows,ncols,Inches(.65),Inches(1.72),Inches(12.05),Inches(4.95)).table
        for r,row in enumerate(data):
            for c in range(ncols):
                cell=shape.cell(r,c); cell.text=row[c] if c<len(row) else ""; cell.fill.solid(); cell.fill.fore_color.rgb=rgb("slate" if r else "slate2"); cell.margin_left=cell.margin_right=Inches(.06)
                for p in cell.text_frame.paragraphs:
                    p.font.name="Aptos"; p.font.size=Pt(10.5 if nrows>7 else 11.5); p.font.bold=(r==0); p.font.color.rgb=rgb("cyan" if r==0 else "ice")
        for c in range(ncols): shape.columns[c].width=Inches(12.05/ncols)
    elif 18 <= idx <= 23:
        pid=sec.select_one(".paper-id"); body=sec.select_one(".paper-body"); year=clean_text(pid.select_one(".year")); pname=clean_text(pid.find("h3")); org=clean_text(pid.select_one(".org")); source=clean_text(pid.select_one(".source")); thesis=clean_text(body.select_one(".thesis")); facts=body.select(".fact"); note=clean_text(body.select_one(".card"))
        add_rect(slide,.65,1.72,4.0,4.95,"slate2","line",True); add_text(slide,year,.95,2.0,3.4,.35,11,"cyan",True,"Courier New"); add_text(slide,pname,.95,2.65,3.35,1.35,25,"white",True,"Aptos Display"); add_text(slide,org,.95,4.35,3.3,.65,13,"muted"); add_text(slide,source,.95,5.75,3.35,.6,9.5,"muted",False,"Courier New")
        add_text(slide,thesis,5.05,1.8,7.45,1.0,19,"white",True)
        y=3.0
        for f in facts[:2]:
            add_rich(slide,clean_text(f.find("b")),clean_text(f.find("p")),5.05,y,3.55,1.55); y2=y
            y=3.0
            if f==facts[0]: pass
        if len(facts)>=2: add_rich(slide,clean_text(facts[1].find("b")),clean_text(facts[1].find("p")),8.85,3.0,3.65,1.55)
        add_rich(slide,"研究判断",note,5.05,4.85,7.45,1.45,"amber")
    elif idx == 26:
        phases=[("2025–2027","受约束闭环","代码、数学、芯片、实验设计；prompt/workflow/code 级改进；人工设定目标与预算。","cyan"),("2027–2030","跨任务迁移","共享 memory 与 skills；改进策略迁移到相邻任务；自动构建 evaluator。","amber"),("2030+","元级递归","改写训练算法与架构，开放式生成新目标；需要更强形式化保证与治理。","risk")]
        for i,(d,t,b,c) in enumerate(phases):
            x=.7+i*4.15; add_rect(slide,x,1.8,3.75,4.85,"slate",c,True); add_text(slide,d,x+.22,2.05,3.1,.3,11,"muted",True,"Courier New"); add_text(slide,t,x+.22,2.7,3.2,.6,22,"white",True); add_text(slide,b,x+.22,3.65,3.25,1.65,14,"ice")
    else:
        # Generic high-density analytic slide: select semantic card-like blocks or paragraphs.
        cards = sec.select(".card")
        if cards:
            n=min(len(cards),6); cols=3 if n>=3 else 2; rows=(n+cols-1)//cols; cw=12.0/cols-.18; ch=4.95/rows-.18
            for i,c in enumerate(cards[:n]):
                x=.65+(i%cols)*(12.0/cols); y=1.72+(i//cols)*(4.95/rows); ct=clean_text(c.find("h3")) or clean_text(c.find(["b","strong"])) or f"要点 {i+1}"; cb=clean_text(c); cb=cb[len(ct):].strip() if cb.startswith(ct) else cb
                add_rich(slide,ct,cb,x,y,cw,ch)
        else:
            paras=[]
            for tag in sec.find_all(["p","li"], recursive=True):
                t=clean_text(tag)
                if t and t not in paras and len(t)>8: paras.append(t)
            if not paras: paras=[t for t in texts if len(t)>10]
            for i,t in enumerate(paras[:6]):
                col=i%2; row=i//2; add_rich(slide,f"要点 {i+1}",t,.7+col*6.08,1.75+row*1.62,5.75,1.38)

# metadata
prs.core_properties.title = "RSI 递归自我提升：硅谷前沿与投资研究"
prs.core_properties.subject = "面向VC与AI行业从业者的30页行业投资研究"
prs.core_properties.author = "DeepSeek Harness"
prs.core_properties.keywords = "RSI, Recursive Self-Improvement, Agentic RSI, Synthetic Self-Training, VC"
prs.save(OUT)
print(OUT)
print("slides", len(prs.slides))
